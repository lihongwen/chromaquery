"""
文件解析器模块
支持多种文档格式的文本提取功能
"""

import logging
import tempfile
import os
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from enum import Enum
import mimetypes

logger = logging.getLogger(__name__)


class FileFormat(Enum):
    """支持的文件格式"""
    TXT = "txt"
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    PPTX = "pptx"
    PPT = "ppt"
    MARKDOWN = "md"
    RTF = "rtf"
    XLSX = "xlsx"
    XLS = "xls"
    CSV = "csv"


@dataclass
class ParseResult:
    """文件解析结果"""
    content: str
    metadata: Dict[str, Any]
    file_format: FileFormat
    success: bool
    error_message: Optional[str] = None
    
    # 表格特有字段
    is_table: bool = False
    table_data: Optional[List[Dict[str, Any]]] = None
    column_analysis: Optional[Dict[str, str]] = None  # 列名 -> 类型(metadata/content)


@dataclass
class TableAnalysisResult:
    """表格分析结果"""
    metadata_columns: List[str]
    content_columns: List[str]
    analysis_reasoning: str


class BaseFileParser(ABC):
    """文件解析器基类"""
    
    def __init__(self):
        self.supported_formats: List[FileFormat] = []
        self.max_file_size = 150 * 1024 * 1024  # 150MB
    
    @abstractmethod
    def can_parse(self, file_format: FileFormat) -> bool:
        """检查是否支持解析指定格式"""
        pass
    
    @abstractmethod
    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析文件内容"""
        pass
    
    def _get_file_format(self, filename: str) -> Optional[FileFormat]:
        """根据文件名获取文件格式"""
        if not filename:
            return None
            
        extension = filename.lower().split('.')[-1] if '.' in filename else ''
        
        format_mapping = {
            'txt': FileFormat.TXT,
            'pdf': FileFormat.PDF,
            'docx': FileFormat.DOCX,
            'doc': FileFormat.DOC,
            'pptx': FileFormat.PPTX,
            'ppt': FileFormat.PPT,
            'md': FileFormat.MARKDOWN,
            'markdown': FileFormat.MARKDOWN,
            'rtf': FileFormat.RTF,
            'xlsx': FileFormat.XLSX,
            'xls': FileFormat.XLS,
            'csv': FileFormat.CSV,
        }
        
        return format_mapping.get(extension)
    
    def _validate_file_size(self, file_content: bytes) -> bool:
        """验证文件大小"""
        return len(file_content) <= self.max_file_size
    
    def _create_temp_file(self, file_content: bytes, suffix: str = '') -> str:
        """创建临时文件"""
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(file_content)
            return temp_file.name
    
    def _cleanup_temp_file(self, temp_path: str):
        """清理临时文件"""
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
        except Exception as e:
            logger.warning(f"清理临时文件失败: {e}")


class TextFileParser(BaseFileParser):
    """文本文件解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.TXT]
    
    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats
    
    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析文本文件"""
        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.TXT,
                    success=False,
                    error_message="文件大小超过限制"
                )
            
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
            content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    content = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.TXT,
                    success=False,
                    error_message="无法解码文件内容，不支持的字符编码"
                )
            
            metadata = {
                "file_size": len(file_content),
                "encoding": used_encoding,
                "line_count": len(content.splitlines()),
                "char_count": len(content)
            }
            
            return ParseResult(
                content=content,
                metadata=metadata,
                file_format=FileFormat.TXT,
                success=True
            )
            
        except Exception as e:
            logger.error(f"解析文本文件失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.TXT,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )


class PDFFileParser(BaseFileParser):
    """PDF文件解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.PDF]
    
    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats
    
    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析PDF文件"""
        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.PDF,
                    success=False,
                    error_message="文件大小超过限制"
                )
            
            # 优先使用pdfplumber，回退到PyPDF2
            content = ""
            metadata = {}
            
            try:
                import pdfplumber
                temp_path = self._create_temp_file(file_content, '.pdf')
                
                try:
                    with pdfplumber.open(temp_path) as pdf:
                        pages_text = []
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                pages_text.append(page_text)
                        
                        content = '\n\n'.join(pages_text)
                        metadata = {
                            "page_count": len(pdf.pages),
                            "file_size": len(file_content),
                            "parser": "pdfplumber"
                        }
                finally:
                    self._cleanup_temp_file(temp_path)
                    
            except ImportError:
                # 回退到PyPDF2
                try:
                    import PyPDF2
                    import io
                    
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    pages_text = []
                    
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            pages_text.append(page_text)
                    
                    content = '\n\n'.join(pages_text)
                    metadata = {
                        "page_count": len(pdf_reader.pages),
                        "file_size": len(file_content),
                        "parser": "PyPDF2"
                    }
                    
                except ImportError:
                    return ParseResult(
                        content="",
                        metadata={},
                        file_format=FileFormat.PDF,
                        success=False,
                        error_message="缺少PDF解析库，请安装pdfplumber或PyPDF2"
                    )
            
            if not content.strip():
                return ParseResult(
                    content="",
                    metadata=metadata,
                    file_format=FileFormat.PDF,
                    success=False,
                    error_message="PDF文件中未提取到文本内容"
                )
            
            return ParseResult(
                content=content,
                metadata=metadata,
                file_format=FileFormat.PDF,
                success=True
            )
            
        except Exception as e:
            logger.error(f"解析PDF文件失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.PDF,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )


class WordFileParser(BaseFileParser):
    """Word文档解析器"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.DOCX, FileFormat.DOC]
    
    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats
    
    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析Word文档"""
        file_format = self._get_file_format(filename)
        
        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="文件大小超过限制"
                )
            
            if file_format == FileFormat.DOCX:
                return self._parse_docx(file_content)
            elif file_format == FileFormat.DOC:
                return self._parse_doc(file_content)
            else:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="不支持的Word文档格式"
                )
                
        except Exception as e:
            logger.error(f"解析Word文档失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=file_format,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )
    
    def _parse_docx(self, file_content: bytes) -> ParseResult:
        """解析DOCX文件"""
        try:
            import docx
            import io
            
            doc = docx.Document(io.BytesIO(file_content))
            paragraphs = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            content = '\n\n'.join(paragraphs)
            
            metadata = {
                "paragraph_count": len(paragraphs),
                "file_size": len(file_content),
                "parser": "python-docx"
            }
            
            return ParseResult(
                content=content,
                metadata=metadata,
                file_format=FileFormat.DOCX,
                success=True
            )
            
        except ImportError:
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.DOCX,
                success=False,
                error_message="缺少python-docx库，请安装python-docx"
            )
    
    def _parse_doc(self, file_content: bytes) -> ParseResult:
        """解析DOC文件（老格式）"""
        # DOC格式解析较复杂，这里提供基础实现
        return ParseResult(
            content="",
            metadata={},
            file_format=FileFormat.DOC,
            success=False,
            error_message="暂不支持.doc格式，请转换为.docx格式"
        )


class PowerPointFileParser(BaseFileParser):
    """PowerPoint文档解析器"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.PPTX, FileFormat.PPT]

    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats

    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析PowerPoint文档"""
        file_format = self._get_file_format(filename)

        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="文件大小超过限制"
                )

            if file_format == FileFormat.PPTX:
                return self._parse_pptx(file_content)
            elif file_format == FileFormat.PPT:
                return self._parse_ppt(file_content)
            else:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="不支持的PowerPoint文档格式"
                )

        except Exception as e:
            logger.error(f"解析PowerPoint文档失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=file_format,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )

    def _parse_pptx(self, file_content: bytes) -> ParseResult:
        """解析PPTX文件"""
        try:
            from pptx import Presentation
            import io

            prs = Presentation(io.BytesIO(file_content))
            slides_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"=== 幻灯片 {slide_num} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                if len(slide_content) > 1:  # 除了标题外还有内容
                    slides_text.append('\n'.join(slide_content))

            content = '\n\n'.join(slides_text)

            metadata = {
                "slide_count": len(prs.slides),
                "file_size": len(file_content),
                "parser": "python-pptx"
            }

            return ParseResult(
                content=content,
                metadata=metadata,
                file_format=FileFormat.PPTX,
                success=True
            )

        except ImportError:
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.PPTX,
                success=False,
                error_message="缺少python-pptx库，请安装python-pptx"
            )

    def _parse_ppt(self, file_content: bytes) -> ParseResult:
        """解析PPT文件（老格式）"""
        return ParseResult(
            content="",
            metadata={},
            file_format=FileFormat.PPT,
            success=False,
            error_message="暂不支持.ppt格式，请转换为.pptx格式"
        )


class MarkdownFileParser(BaseFileParser):
    """Markdown文件解析器"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.MARKDOWN]

    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats

    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析Markdown文件"""
        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.MARKDOWN,
                    success=False,
                    error_message="文件大小超过限制"
                )

            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
            content = None
            used_encoding = None

            for encoding in encodings:
                try:
                    content = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.MARKDOWN,
                    success=False,
                    error_message="无法解码文件内容，不支持的字符编码"
                )

            # 可选：转换Markdown为纯文本
            try:
                import markdown
                from bs4 import BeautifulSoup

                # 转换Markdown为HTML
                html = markdown.markdown(content)
                # 提取纯文本
                soup = BeautifulSoup(html, 'html.parser')
                plain_text = soup.get_text()

                metadata = {
                    "file_size": len(file_content),
                    "encoding": used_encoding,
                    "line_count": len(content.splitlines()),
                    "parser": "markdown+beautifulsoup",
                    "original_markdown": True
                }

                return ParseResult(
                    content=plain_text,
                    metadata=metadata,
                    file_format=FileFormat.MARKDOWN,
                    success=True
                )

            except ImportError:
                # 如果没有markdown库，直接返回原始内容
                metadata = {
                    "file_size": len(file_content),
                    "encoding": used_encoding,
                    "line_count": len(content.splitlines()),
                    "parser": "raw_text",
                    "original_markdown": True
                }

                return ParseResult(
                    content=content,
                    metadata=metadata,
                    file_format=FileFormat.MARKDOWN,
                    success=True
                )

        except Exception as e:
            logger.error(f"解析Markdown文件失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.MARKDOWN,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )


class RTFFileParser(BaseFileParser):
    """RTF文件解析器"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.RTF]

    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats

    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析RTF文件"""
        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.RTF,
                    success=False,
                    error_message="文件大小超过限制"
                )

            try:
                from striprtf.striprtf import rtf_to_text

                # 尝试解码RTF内容
                try:
                    rtf_content = file_content.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        rtf_content = file_content.decode('latin-1')
                    except UnicodeDecodeError:
                        return ParseResult(
                            content="",
                            metadata={},
                            file_format=FileFormat.RTF,
                            success=False,
                            error_message="无法解码RTF文件内容"
                        )

                # 提取纯文本
                content = rtf_to_text(rtf_content)

                metadata = {
                    "file_size": len(file_content),
                    "parser": "striprtf"
                }

                return ParseResult(
                    content=content,
                    metadata=metadata,
                    file_format=FileFormat.RTF,
                    success=True
                )

            except ImportError:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.RTF,
                    success=False,
                    error_message="缺少striprtf库，请安装striprtf"
                )

        except Exception as e:
            logger.error(f"解析RTF文件失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.RTF,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )


class TableFileParser(BaseFileParser):
    """表格文件解析器（Excel和CSV）"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [FileFormat.XLSX, FileFormat.XLS, FileFormat.CSV]

    def can_parse(self, file_format: FileFormat) -> bool:
        return file_format in self.supported_formats

    def parse(self, file_content: bytes, filename: str) -> ParseResult:
        """解析表格文件"""
        file_format = self._get_file_format(filename)

        try:
            if not self._validate_file_size(file_content):
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="文件大小超过限制"
                )

            if file_format in [FileFormat.XLSX, FileFormat.XLS]:
                return self._parse_excel(file_content, file_format)
            elif file_format == FileFormat.CSV:
                return self._parse_csv(file_content)
            else:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="不支持的表格文件格式"
                )

        except Exception as e:
            logger.error(f"解析表格文件失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=file_format,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )

    def _parse_excel(self, file_content: bytes, file_format: FileFormat) -> ParseResult:
        """解析Excel文件"""
        try:
            import pandas as pd
            import io
            import numpy as np

            # 读取Excel文件
            df = pd.read_excel(io.BytesIO(file_content), engine='openpyxl' if file_format == FileFormat.XLSX else 'xlrd')

            # 检查是否为空
            if df.empty:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="Excel文件为空"
                )

            # 清理数据
            df = self._clean_excel_data(df)

            # 再次检查清理后是否为空
            if df.empty:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message="Excel文件清理后为空"
                )

            # 转换为字典列表
            table_data = df.to_dict('records')

            # 智能分析列类型
            column_analysis = self._analyze_table_columns(df)

            # 生成文档内容（每行一个文档）
            documents = self._convert_table_to_documents(table_data, column_analysis)

            metadata = {
                "file_size": len(file_content),
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": ", ".join(str(col) for col in df.columns),  # 转换为字符串
                "parser": "pandas"
            }

            return ParseResult(
                content="\n\n".join(documents),
                metadata=metadata,
                file_format=file_format,
                success=True,
                is_table=True,
                table_data=table_data,
                column_analysis=column_analysis
            )

        except ImportError:
            return ParseResult(
                content="",
                metadata={},
                file_format=file_format,
                success=False,
                error_message="缺少pandas或openpyxl库，请安装相关依赖"
            )
        except Exception as e:
            logger.error(f"Excel解析失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=file_format,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )

    def _clean_excel_data(self, df):
        """清理Excel数据"""
        import pandas as pd
        import numpy as np

        # 1. 删除完全为空的行
        df = df.dropna(how='all')

        # 2. 删除完全为空的列
        df = df.dropna(axis=1, how='all')

        # 3. 清理列名
        # 移除列名中的换行符和多余空格
        df.columns = [str(col).replace('\n', ' ').replace('\r', ' ').strip() for col in df.columns]

        # 4. 处理Unnamed列 - 如果有相邻的有意义列名，则合并
        new_columns = []
        for i, col in enumerate(df.columns):
            if str(col).startswith('Unnamed'):
                # 查找前一个有意义的列名
                if i > 0 and not str(df.columns[i-1]).startswith('Unnamed'):
                    # 如果这个Unnamed列有数据，则重命名为前一列的扩展
                    if not df[col].isna().all():
                        new_columns.append(f"{df.columns[i-1]}_详细")
                    else:
                        new_columns.append(col)  # 保持原名，后续删除
                else:
                    new_columns.append(col)
            else:
                new_columns.append(col)

        df.columns = new_columns

        # 5. 删除以"Unnamed"开头的空列
        unnamed_cols = [col for col in df.columns if str(col).startswith('Unnamed') and df[col].isna().all()]
        df = df.drop(columns=unnamed_cols)

        # 6. 过滤掉主要内容为空的行（保留至少有3个非空值的行）
        min_non_null = max(1, len(df.columns) // 3)  # 至少1/3的列有值
        df = df.dropna(thresh=min_non_null)

        # 7. 重置索引
        df = df.reset_index(drop=True)

        return df

    def _parse_csv(self, file_content: bytes) -> ParseResult:
        """解析CSV文件"""
        try:
            import pandas as pd
            import io

            # 尝试不同的编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
            df = None
            used_encoding = None

            for encoding in encodings:
                try:
                    content_str = file_content.decode(encoding)
                    df = pd.read_csv(io.StringIO(content_str))
                    used_encoding = encoding
                    break
                except (UnicodeDecodeError, pd.errors.EmptyDataError):
                    continue

            if df is None or df.empty:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=FileFormat.CSV,
                    success=False,
                    error_message="无法解析CSV文件或文件为空"
                )

            # 转换为字典列表
            table_data = df.to_dict('records')

            # 智能分析列类型
            column_analysis = self._analyze_table_columns(df)

            # 生成文档内容（每行一个文档）
            documents = self._convert_table_to_documents(table_data, column_analysis)

            metadata = {
                "file_size": len(file_content),
                "encoding": used_encoding,
                "row_count": len(df),
                "column_count": len(df.columns),
                "columns": ", ".join(str(col) for col in df.columns),  # 转换为字符串
                "parser": "pandas"
            }

            return ParseResult(
                content="\n\n".join(documents),
                metadata=metadata,
                file_format=FileFormat.CSV,
                success=True,
                is_table=True,
                table_data=table_data,
                column_analysis=column_analysis
            )

        except ImportError:
            return ParseResult(
                content="",
                metadata={},
                file_format=FileFormat.CSV,
                success=False,
                error_message="缺少pandas库，请安装pandas"
            )

    def _analyze_table_columns(self, df) -> Dict[str, str]:
        """使用LLM智能分析表格列类型（基于标题行和前10行数据样本）"""
        try:
            logger.info("开始LLM智能列分析")

            # 调用LLM分析（传入标题行和前10行数据）
            analysis_result = self._call_llm_for_column_analysis_with_sample_data(df)
            if analysis_result:
                logger.info(f"LLM列分析成功: {analysis_result}")
                return analysis_result
            else:
                logger.warning("LLM列分析失败，使用简单规则分析")
                return self._simple_column_analysis(df)

        except Exception as e:
            logger.warning(f"智能列分析失败，使用简单规则: {e}")
            return self._simple_column_analysis(df)

    def _call_llm_for_column_analysis_with_sample_data(self, df) -> Optional[Dict[str, str]]:
        """调用LLM分析表格列（基于标题行和前10行数据样本）"""
        try:
            from llm_client import get_llm_client

            llm_client = get_llm_client()
            if not llm_client:
                logger.warning("LLM客户端不可用")
                return None

            # 构建分析提示（包含标题行和前10行数据样本）
            prompt = self._build_column_analysis_prompt_with_sample_data(df)

            # 调用LLM（同步方式）
            import asyncio
            messages = [{"role": "user", "content": prompt}]

            # 检查是否已有运行中的事件循环
            try:
                # 尝试获取当前事件循环
                loop = asyncio.get_running_loop()
                # 如果有运行中的循环，使用 asyncio.create_task 在当前循环中运行
                import concurrent.futures

                async def collect_response():
                    response_parts = []
                    async for chunk in llm_client.stream_chat(messages, temperature=0.1, max_tokens=1500):
                        if chunk.get('content'):
                            response_parts.append(chunk['content'])
                    return ''.join(response_parts)

                # 在线程池中运行异步函数
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(asyncio.run, collect_response())
                    response = future.result()

            except RuntimeError:
                # 没有运行中的事件循环，创建新的
                async def collect_response():
                    response_parts = []
                    async for chunk in llm_client.stream_chat(messages, temperature=0.1, max_tokens=1500):
                        if chunk.get('content'):
                            response_parts.append(chunk['content'])
                    return ''.join(response_parts)

                response = asyncio.run(collect_response())

            logger.info(f"LLM分析响应: {response[:200]}...")  # 记录前200字符

            # 解析响应
            return self._parse_llm_column_response(response)

        except Exception as e:
            logger.error(f"LLM列分析失败: {e}")
            return None

    def _call_llm_for_column_analysis(self, columns_info: Dict) -> Optional[Dict[str, str]]:
        """调用LLM分析表格列（旧版本，保留兼容性）"""
        try:
            from llm_client import get_llm_client

            llm_client = get_llm_client()
            if not llm_client:
                return None

            # 构建分析提示
            prompt = self._build_column_analysis_prompt(columns_info)

            # 调用LLM（同步方式）
            import asyncio
            messages = [{"role": "user", "content": prompt}]

            # 创建事件循环并运行异步方法
            try:
                loop = asyncio.get_event_loop()
            except RuntimeError:
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)

            # 收集流式响应
            response_parts = []
            async def collect_response():
                async for chunk in llm_client.stream_chat(messages, temperature=0.1, max_tokens=1000):
                    if chunk.get('content'):
                        response_parts.append(chunk['content'])

            loop.run_until_complete(collect_response())
            response = ''.join(response_parts)

            # 解析响应
            return self._parse_llm_column_response(response)

        except Exception as e:
            logger.error(f"LLM列分析失败: {e}")
            return None

    def _build_column_analysis_prompt(self, columns_info: Dict) -> str:
        """构建列分析提示（使用完整数据样本）"""
        # 从columns_info中获取DataFrame（需要修改调用方式）
        # 这里先用原有逻辑，稍后会修改整个流程
        prompt = """请分析以下表格的列，判断每列应该作为元数据(metadata)还是内容(content)存储到向量数据库中。

分析规则：
- **content列**：包含丰富文本信息，适合语义搜索的列
  * 产品名称、标题、描述、内容、评论、摘要等
  * 通常包含多个词汇，具有语义含义
  * 用户会基于这些内容进行搜索查询

- **metadata列**：结构化信息，用于筛选和标识的列
  * ID、编号、代码、分类、标签、价格、数量、时间、状态等
  * 通常是标识符、数值、分类或时间信息
  * 用于过滤和组织数据，而非语义搜索

表格列信息：
"""

        for col, info in columns_info.items():
            prompt += f"\n列名: {col}\n"
            prompt += f"数据类型: {info['dtype']}\n"
            prompt += f"示例值: {info['sample_values']}\n"
            prompt += f"唯一值数量: {info['unique_count']}\n"
            prompt += "---\n"

        prompt += """
请以JSON格式返回分析结果，格式如下：
{
    "column_name1": "metadata",
    "column_name2": "content",
    ...
}

只返回JSON，不要其他解释。"""

        return prompt

    def _build_column_analysis_prompt_with_data(self, df) -> str:
        """构建列分析提示（包含完整数据样本）"""
        # 限制数据行数，避免token过多
        sample_rows = min(15, len(df))
        df_sample = df.head(sample_rows)

        prompt = f"""请分析以下表格数据，判断每列应该作为元数据(metadata)还是内容(content)存储到向量数据库中。

**分析规则：**

**content列** - 适合语义搜索的文本内容：
- 产品名称、标题、描述、内容、评论、摘要等
- 包含丰富的文本信息，用户会基于这些内容进行搜索
- 通常包含多个词汇，具有语义含义
- 例如：商品描述、文章标题、用户评论等

**metadata列** - 结构化标识信息：
- ID、编号、代码、分类、价格、数量、时间、状态等
- 用于筛选、过滤和组织数据
- 通常是标识符、数值、分类标签或时间信息
- 例如：商品ID、价格、分类、创建时间等

**表格数据样本（共{len(df)}行，显示前{sample_rows}行）：**

"""

        # 将DataFrame转换为易读的表格格式
        # 表头
        headers = list(df_sample.columns)
        prompt += "| " + " | ".join(headers) + " |\n"
        prompt += "|" + "|".join([" --- " for _ in headers]) + "|\n"

        # 数据行
        for idx, row in df_sample.iterrows():
            row_data = []
            for col in headers:
                value = str(row[col]) if row[col] is not None else ""
                # 限制单元格长度，避免过长
                if len(value) > 50:
                    value = value[:47] + "..."
                row_data.append(value)
            prompt += "| " + " | ".join(row_data) + " |\n"

        prompt += f"""

**请基于以上完整的数据样本进行分析，考虑：**
1. 每列的实际内容和数据特征
2. 哪些列包含丰富的文本信息，适合语义搜索
3. 哪些列是结构化信息，适合作为筛选条件

请以JSON格式返回分析结果：
{{
    "列名1": "content",
    "列名2": "metadata",
    ...
}}

只返回JSON，不要其他解释。"""

        return prompt

    def _build_column_analysis_prompt_with_sample_data(self, df) -> str:
        """构建列分析提示（基于标题行和前10行数据样本）"""
        # 获取前10行数据作为样本
        sample_size = min(10, len(df))
        sample_df = df.head(sample_size)

        # 构建数据样本字符串
        data_sample = "**表格数据样本：**\n"
        data_sample += "**标题行（列名）：** " + " | ".join(str(col) for col in df.columns) + "\n\n"
        data_sample += "**前{}行数据：**\n".format(sample_size)

        for idx, row in sample_df.iterrows():
            row_data = []
            for col, value in zip(df.columns, row.values):
                if value is not None and str(value).strip():
                    row_data.append(f"{col}: {str(value).strip()}")
                else:
                    row_data.append(f"{col}: [空]")
            data_sample += f"第{idx+1}行: {' | '.join(row_data)}\n"

        prompt = f"""请分析以下Excel表格的列，判断每列应该作为元数据(metadata)还是内容(content)存储到向量数据库中。

**分析规则：**

**content列** - 包含丰富文本信息，适合语义搜索：
- 产品名称、标题、描述、内容、评论、摘要、说明等
- 通常包含多个词汇，具有语义含义
- 用户会基于这些内容进行搜索查询
- 例如：产品描述、文章标题、用户评论等

**metadata列** - 结构化信息，用于筛选和标识：
- ID、编号、代码、分类、标签、价格、数量、时间、状态等
- 通常是标识符、数值、分类或时间信息
- 用于过滤和组织数据，而非语义搜索
- 例如：产品ID、价格、分类、创建时间等

{data_sample}

**请基于以上标题行和数据样本进行分析，考虑：**
1. 每列的实际内容和数据特征
2. 哪些列包含丰富的文本信息，适合语义搜索
3. 哪些列是结构化信息，适合作为筛选条件

**请以JSON格式返回分析结果：**
{{
    "列名1": "content",
    "列名2": "metadata",
    ...
}}

只返回JSON，不要其他解释。"""

        return prompt

    def _parse_llm_column_response(self, response: str) -> Optional[Dict[str, str]]:
        """解析LLM的列分析响应"""
        try:
            import json
            import re

            # 提取JSON部分
            json_match = re.search(r'\{[^}]+\}', response, re.DOTALL)
            if json_match:
                json_str = json_match.group()
                return json.loads(json_str)

            return None

        except Exception as e:
            logger.error(f"解析LLM响应失败: {e}")
            return None

    def _simple_column_analysis(self, df) -> Dict[str, str]:
        """简单的列类型分析（回退方案）"""
        column_analysis = {}

        for col in df.columns:
            col_lower = col.lower()

            # 优先识别content列（包含丰富文本信息的列）
            if any(keyword in col_lower for keyword in [
                'name', 'title', 'content', 'description', 'text', 'comment', 'summary',
                '名称', '标题', '内容', '描述', '评论', '摘要', '详情', '说明'
            ]):
                column_analysis[col] = 'content'
            # 明确的metadata列
            elif any(keyword in col_lower for keyword in [
                'id', '编号', 'number', 'code', '代码', 'price', '价格', 'amount', '金额',
                'quantity', '数量', '库存', 'stock', 'count', '计数'
            ]):
                column_analysis[col] = 'metadata'
            elif any(keyword in col_lower for keyword in [
                'time', 'date', '时间', '日期', 'created', 'updated', '创建', '更新'
            ]):
                column_analysis[col] = 'metadata'
            elif any(keyword in col_lower for keyword in [
                'status', 'state', '状态', 'type', 'category', '分类', '类型', 'tag', '标签'
            ]):
                column_analysis[col] = 'metadata'
            else:
                # 根据数据内容特征判断
                if df[col].dtype == 'object':
                    # 计算平均文本长度
                    avg_length = df[col].astype(str).str.len().mean()
                    # 计算唯一值比例
                    unique_ratio = df[col].nunique() / len(df)

                    # 如果平均长度较长且唯一值比例较高，可能是content
                    if avg_length > 15 and unique_ratio > 0.5:
                        column_analysis[col] = 'content'
                    else:
                        column_analysis[col] = 'metadata'
                else:
                    # 数值类型通常是metadata
                    column_analysis[col] = 'metadata'

        return column_analysis

    def _convert_table_to_documents(self, table_data: List[Dict], column_analysis: Dict[str, str]) -> List[str]:
        """将表格数据转换为文档列表"""
        documents = []

        content_columns = [col for col, type_ in column_analysis.items() if type_ == 'content']

        for row_idx, row in enumerate(table_data):
            # 提取内容列
            content_parts = []
            for col in content_columns:
                if col in row and row[col] is not None and str(row[col]).strip():
                    content_parts.append(f"{col}: {str(row[col]).strip()}")

            if content_parts:
                document = " | ".join(content_parts)
                documents.append(document)
            else:
                # 如果没有内容列，将所有列作为内容
                all_parts = []
                for col, value in row.items():
                    if value is not None and str(value).strip():
                        all_parts.append(f"{col}: {str(value).strip()}")
                if all_parts:
                    documents.append(" | ".join(all_parts))

        return documents


class FileParserManager:
    """文件解析器管理器"""

    def __init__(self):
        self.parsers = [
            TextFileParser(),
            PDFFileParser(),
            WordFileParser(),
            PowerPointFileParser(),
            MarkdownFileParser(),
            RTFFileParser(),
            TableFileParser(),
        ]

        # 构建格式到解析器的映射
        self.format_parser_map = {}
        for parser in self.parsers:
            for format_ in parser.supported_formats:
                self.format_parser_map[format_] = parser

    def get_supported_formats(self) -> List[FileFormat]:
        """获取所有支持的文件格式"""
        return list(self.format_parser_map.keys())

    def get_supported_extensions(self) -> List[str]:
        """获取所有支持的文件扩展名"""
        extensions = []
        for format_ in self.get_supported_formats():
            extensions.append(f".{format_.value}")
        return extensions

    def can_parse(self, filename: str) -> bool:
        """检查是否支持解析指定文件"""
        file_format = self._get_file_format(filename)
        return file_format in self.format_parser_map

    def parse_file(self, file_content: bytes, filename: str) -> ParseResult:
        """解析文件"""
        try:
            # 获取文件格式
            file_format = self._get_file_format(filename)
            if not file_format:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=None,
                    success=False,
                    error_message="无法识别文件格式"
                )

            # 获取对应的解析器
            parser = self.format_parser_map.get(file_format)
            if not parser:
                return ParseResult(
                    content="",
                    metadata={},
                    file_format=file_format,
                    success=False,
                    error_message=f"不支持的文件格式: {file_format.value}"
                )

            # 执行解析
            return parser.parse(file_content, filename)

        except Exception as e:
            logger.error(f"文件解析失败: {e}")
            return ParseResult(
                content="",
                metadata={},
                file_format=None,
                success=False,
                error_message=f"解析失败: {str(e)}"
            )

    def _get_file_format(self, filename: str) -> Optional[FileFormat]:
        """根据文件名获取文件格式"""
        if not filename:
            return None

        extension = filename.lower().split('.')[-1] if '.' in filename else ''

        format_mapping = {
            'txt': FileFormat.TXT,
            'pdf': FileFormat.PDF,
            'docx': FileFormat.DOCX,
            'doc': FileFormat.DOC,
            'pptx': FileFormat.PPTX,
            'ppt': FileFormat.PPT,
            'md': FileFormat.MARKDOWN,
            'markdown': FileFormat.MARKDOWN,
            'rtf': FileFormat.RTF,
            'xlsx': FileFormat.XLSX,
            'xls': FileFormat.XLS,
            'csv': FileFormat.CSV,
        }

        return format_mapping.get(extension)


# 全局文件解析器管理器实例
file_parser_manager = FileParserManager()
